from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

ASSETS = r"c:\aidevme-blog\aidevme-blog-content\draft-articles\60-mda-webresource-js-vs-ts\assets"
OUT    = r"c:\aidevme-blog\aidevme-blog-content\draft-articles\60-mda-webresource-js-vs-ts\presentation.pptx"

# 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

BG      = RGBColor(0x0A, 0x0E, 0x1A)   # near-black
CYAN    = RGBColor(0x00, 0xD4, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x88, 0x99, 0xAA)
ACCENT  = RGBColor(0x00, 0xFF, 0xCC)

def img(name):
    return os.path.join(ASSETS, name)

def set_bg(slide, prs):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_label(slide, text, left, top, width, height,
              size=11, bold=False, color=GREY, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

# ── slide helpers ─────────────────────────────────────────────────────────────

def title_slide(prs, title, subtitle, img_file, img_alpha=True):
    """Full-bleed image on right half, text on left."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, prs)

    # image — right 60 %
    img_l = Inches(5.0)
    slide.shapes.add_picture(img(img_file), img_l, Inches(0), Inches(8.33), H)

    # dark overlay strip on left
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(5.4), H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()

    # aidevme label
    add_label(slide, "aidevme.com", Inches(0.4), Inches(0.3), Inches(4.6), Inches(0.4),
              size=10, color=CYAN)

    # title
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(4.6), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # subtitle
    add_label(slide, subtitle, Inches(0.4), Inches(5.0), Inches(4.6), Inches(1.8),
              size=14, color=GREY)

    # bottom accent line
    ln = slide.shapes.add_shape(1, Inches(0.4), Inches(6.9), Inches(2.0), Emu(36000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = CYAN
    ln.line.fill.background()

    return slide


def full_image_slide(prs, heading, img_file, note=None):
    """Image takes up most of the slide, heading at top."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)

    # heading bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0D, 0x14, 0x26)
    bar.line.fill.background()

    add_label(slide, heading, Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.55),
              size=18, bold=True, color=WHITE)

    # image centered below header
    img_top = Inches(0.8)
    img_h   = Inches(6.1)
    # calculate width to preserve ratio
    from PIL import Image as PILImage
    path = img(img_file)
    with PILImage.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    img_w = img_h * ratio
    if img_w > W:
        img_w = W
        img_h = img_w / ratio
    img_left = (W - img_w) / 2
    slide.shapes.add_picture(path, img_left, img_top, img_w, img_h)

    if note:
        add_label(slide, note, Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.35),
                  size=9, color=GREY, align=PP_ALIGN.CENTER)

    # slide number dot
    add_label(slide, "aidevme.com", W - Inches(1.5), Inches(7.1), Inches(1.4), Inches(0.35),
              size=9, color=CYAN, align=PP_ALIGN.RIGHT)
    return slide


def cta_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)

    # glow rectangle
    box = slide.shapes.add_shape(1, Inches(1.5), Inches(1.0), Inches(10.33), Inches(5.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x0D, 0x1E, 0x30)
    box.line.color.rgb = CYAN
    box.line.width = Emu(25400)

    add_label(slide, "Read the full guide →", Inches(2.0), Inches(1.5), Inches(9.0), Inches(0.8),
              size=14, color=CYAN, align=PP_ALIGN.CENTER)

    add_label(slide,
              "JavaScript vs TypeScript Web Resources\nin Model-Driven Apps: Complete Guide (2026)",
              Inches(2.0), Inches(2.2), Inches(9.0), Inches(1.4),
              size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_label(slide,
              "aidevme.com/javascript-vs-typescript-web-resources-in-model-driven-apps-complete-guide/",
              Inches(2.0), Inches(3.7), Inches(9.0), Inches(0.5),
              size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    add_label(slide,
              "Like this? Follow for more Power Platform architecture deep dives.",
              Inches(2.0), Inches(4.5), Inches(9.0), Inches(0.5),
              size=13, color=GREY, align=PP_ALIGN.CENTER)

    # accent line
    ln = slide.shapes.add_shape(1, Inches(5.2), Inches(5.3), Inches(2.9), Emu(36000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = CYAN
    ln.line.fill.background()

    add_label(slide, "aidevme.com", Inches(2.0), Inches(6.9), Inches(9.0), Inches(0.4),
              size=10, color=CYAN, align=PP_ALIGN.CENTER)
    return slide

# ── build ─────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# 1 — Title
title_slide(
    prs,
    "JavaScript vs TypeScript\nWeb Resources in\nModel-Driven Apps",
    "Plain JS · Webpack · esbuild · npm Workspaces · CI/CD\nComplete Guide 2026",
    "javascript-vs-typescript-web-resources-model-driven-apps-hero.png"
)

# 2 — Why Now
full_image_slide(
    prs,
    "Why This Question Matters Now",
    "avascript-vs-typescript-web-resources-model-driven-apps-why-this-question-matters-now.png",
    "From legacy CRM scripts to modern TypeScript toolchains"
)

# 3 — The Example
full_image_slide(
    prs,
    "The Example We'll Use: Contact Form",
    "javascript-vs-typescript-web-resources-model-driven-apps-contact-form-example-overview.png",
    "OnChange event → greeting notification + WebApi account lookup + validation"
)

# 4 — Plain JS
full_image_slide(
    prs,
    "Approach 1: Plain JavaScript",
    "javascript-vs-typescript-web-resources-model-driven-apps-approach-plain-javascript.png",
    "Zero build step · No toolchain · Fast to ship · No type safety"
)

# 5 — Webpack
full_image_slide(
    prs,
    "Approach 2: TypeScript + Webpack",
    "javascript-vs-typescript-web-resources-model-driven-apps-approach-typescript-webpack.png",
    "Full type safety · Rich plugin ecosystem · Slower builds · Complex config"
)

# 6 — esbuild
full_image_slide(
    prs,
    "Approach 3: TypeScript + esbuild",
    "javascript-vs-typescript-web-resources-model-driven-apps-approach-typescript-esbuild.png",
    "Sub-100ms builds · Minimal config · Same type safety · CI/CD-friendly"
)

# 7 — When to choose
full_image_slide(
    prs,
    "When to Choose What",
    "javascript-vs-typescript-web-resources-model-driven-apps-when-to-choose-what-decision-guide.png",
    "Raw JS = solo/quick · Webpack = complex team · esbuild = modern default"
)

# 8 — npm workspaces
full_image_slide(
    prs,
    "Multi-Developer Projects: npm Workspace Architecture",
    "javascript-vs-typescript-web-resources-model-driven-apps-multi-developer-npm-workspaces-architecture.png",
    "@aidevme/shared consumed by all form and ribbon packages — no merge conflicts"
)

# 9 — Manual upload
full_image_slide(
    prs,
    "Deployment Scenario 1: Manual Upload via Maker Portal",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-manual-upload-maker-portal.png",
    "Build → Upload to make.powerapps.com → Publish · Best for first-time setup"
)

# 10 — pac CLI
full_image_slide(
    prs,
    "Deployment Scenario 2: Power Platform CLI — Local Workflow",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-pac-cli-local-workflow.png",
    "Edit → npm run build → pac webresource upload → pac solution publish"
)

# 11 — GitHub Actions
full_image_slide(
    prs,
    "Deployment Scenario 3: GitHub Actions CI/CD",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-github-actions-cicd.png",
    "Push → Typecheck → Build → pac auth → pac upload → pac publish · Automated on every merge"
)

# 12 — Azure DevOps
full_image_slide(
    prs,
    "Deployment Scenario 4: Azure DevOps Multi-Stage Pipeline",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-azure-devops-pipelines.png",
    "CI Build → Deploy Test (optional approval) → Deploy Production (required approval)"
)

# 13 — Managed vs Unmanaged
full_image_slide(
    prs,
    "Deployment Scenario 5: Managed vs Unmanaged Solutions",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-managed-vs-unmanaged-solutions.png",
    "Dev = Unmanaged · Test & Production = Managed · Always version-bump on release"
)

# 14 — Decision ladder
full_image_slide(
    prs,
    "Deployment Scenario Decision Guide",
    "javascript-vs-typescript-web-resources-model-driven-apps-deployment-scenario-decision-guide.png",
    "Start simple, automate incrementally — each step builds on the last"
)

# 15 — Conclusion
full_image_slide(
    prs,
    "Conclusion: The Three Paths to Production",
    "javascript-vs-typescript-web-resources-model-driven-apps-conclusion-decision-summary.png",
    "Solo dev = plain JS · Small team = esbuild + pac CLI · Enterprise = workspaces + pipelines + managed solutions"
)

# 16 — CTA
cta_slide(prs)

prs.save(OUT)
print(f"Saved: {OUT}")
